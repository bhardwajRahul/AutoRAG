import glob, os
from pptx import Presentation
from pptx.util import Inches

PNG_DIR = "decks/journey-of-autorag/preview-png"
OUT = "decks/journey-of-autorag-images.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

pngs = sorted(glob.glob(os.path.join(PNG_DIR, "slide-*.png")))
for p in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print(f"wrote {OUT} with {len(pngs)} full-bleed image slides")
